from django.shortcuts import render, redirect
from django.http import HttpResponse
from django.contrib import messages
from django.views.generic import TemplateView, View

from pages.models import TopResearches


from my_requests.models import TranslationCostModel
from orders.models import OrderModel


from .forms import (ConsultationForm, FreeConsultationForm, ProofreadingForm, BaksForm, ScopusesForm,
                    PeerReviewForm, OrganizeConferencesForm, GrantsForm, TranslationForm, PatentsForm, DesignsForm)
from .models import Translation, ResearchPlatformsContext, GetDesign

from setpage.forms import WebinarsForm, TranslationsForm

from setpage.models import ( WebinarsModel, TranslationModel, BakModel, ConferencesModel, CreateConferenceModel, DesignModel, GrantsModel,
                            PatentModel, ScopusModel, ProofModel, WebinarsUrls1)
from setpage.forms import BakForm, ConferencesForm, CreateConferenceForm, DesignForm, GrantForm, PatentForm, ScopusForm, \
    ProofForm


class Research(TemplateView):
    template_name = 'research_platform.html'

    def get_context_data(self, **kwargs):
        context = super(Research, self).get_context_data(**kwargs)
        context['text'] = ResearchPlatformsContext.objects.all()
        context['title'] = 'Research Platform'
        return context


class ResearchIntelligense(TemplateView):
    template_name = 'research_intelligense.html'

    def get_context_data(self, **kwargs):
        context = super(ResearchIntelligense, self).get_context_data(**kwargs)
        context['title'] = 'Research Intelligens'
        return context


class Scientific(TemplateView):
    template_name = 'scientific_papers.html'

    def get_context_data(self, **kwargs):
        context = super(Scientific, self).get_context_data(**kwargs)
        context['title'] = 'Scientific Papers'
        return context


class TopResearchess(TemplateView):
    template_name = 'top_researches.html'

    def get_context_data(self, **kwargs):
        context = super(TopResearchess, self).get_context_data(**kwargs)
        context['top_researches'] = TopResearches.objects.all()
        context['title'] = 'Top Researches'
        return context


class WebinarsView(TemplateView):
    template_name = 'webinars.html'


    def get_context_data(self, **kwargs):
        context = super(WebinarsView, self).get_context_data(**kwargs)
        text = WebinarsModel.objects.first()
        context['webinar'] = WebinarsForm(instance=text)
        context['text'] = text
        context['url1'] = WebinarsUrls1.objects.filter().order_by('-id')[:5]
        context['url2'] = WebinarsUrls1.objects.filter().order_by('-id')[5:10]
        context['title'] = 'Webinars'
        return context


# todo yangilash
class Conferences(View):
    def get(self, request):
        text = ConferencesModel.objects.first()
        con = ConferencesForm(instance=text)
        context = {'form': FreeConsultationForm(), 'con': con, 'text': text, 'title':'Conferencess'}
        return render(request, 'conferences.html', context)

    def post(self, request):
        form = FreeConsultationForm(request.POST, request.FILES)

        if request.user.is_authenticated:
            if form.is_valid():
                if request.POST.get('is_agree') and request.POST.get('is_agree') == 'on':
                    done = form.save()
                    if done:
                        messages.success(request, 'Ваш запрос успешно отправлен, мы скоро свяжемся с вами!')
                        return redirect('requests:conferences')
                else:
                    messages.error(request, 'вы  должны согласиться прежде чем отправить форму ')
                    return render(request, 'conferences.html', {'validated': 'validated', 'form': form})
            else:
                return render(request, 'conferences.html', {'validated': 'validated', 'form': form})
        else:
            messages.error(request, 'Чтобы отправить форму, вы должны сначала войти в систему')
            return redirect('users:login')



class CreateConferences(View):
    def get(self, request):
        text = CreateConferenceModel.objects.first()
        con = CreateConferenceForm(instance=text)
        context = {'form': OrganizeConferencesForm(), 'text': text, 'con': con, 'title':'Create Conferencess'}
        return render(request, 'creat_conference.html', context)

    def post(self, request):
        form = OrganizeConferencesForm(request.POST)
        if form.is_valid():  
            done = form.save()
            if done:
                messages.success(request, 'Ваш запрос успешно отправлен, мы скоро свяжемся с вами!')
                return redirect('requests:creat-conferences')
        return render(request, 'creat_conference.html', {'validated': 'validated', 'form': form})
      

class Design(View):
    def get(self, request):
        obj = DesignModel.objects.first()
        design = DesignForm(instance=obj)
        context = {'form': DesignsForm(), 'text': obj, 'design': design, 'title':'Design'}
        return render(request, 'design.html', context)

    def post(self, request):
        form = DesignsForm(request.POST)
        
        if form.is_valid():
            done = form.save()
            if done:
                messages.success(request, 'Ваш запрос успешно отправлен, мы скоро свяжемся с вами!')
                return redirect('requests:design')
        else:
            return render(request, 'design.html', {'validated': 'validated', 'form': form})
       


class ConsultationView(View):
    def get(self, request):
        form = ConsultationForm()
        return render(request, 'home.html', {'form': form, 'title': 'Consultation'})

    def post(self, request):
        form = ConsultationForm(request.POST)
        
        if form.is_valid():
            
            done = form.save()
            if done:
                messages.success(request, 'Ваш запрос успешно отправлен, мы скоро свяжемся с вами!')
                return redirect('pages:home')
            
        else:
            return render(request, 'home.html', {'validated': 'validated', 'form': form})
    


# class FreeConsultationView(View):
#    def post(self, request):
#         form = FreeConsultationForm(request.POST, request.FILES)
#         if request.user.is_authenticated:
#             if form.is_valid():
#                 if request.POST.get('is_agree') and request.POST.get('is_agree') == 'on':
#                     done = form.save()
#                     if done:
#                         messages.success(request, 'Ваш запрос успешно отправлен, мы скоро свяжемся с вами!')
#                         return redirect('requests:conferences')
#                 else:
#                     messages.error(request, 'вы  должны согласиться прежде чем отправить форму ')
#                     return render(request,'conferences.html', {'validated':'validated', 'form':form})
#             else:
#                 return render(request,'conferences.html', {'validated':'validated', 'form':form})
#         else:
#             messages.error(request, 'Чтобы отправить форму, вы должны сначала войти в систему')
#             return redirect('users:login')


class Proofreading(View):
    def get(self, request):
        data = request.GET.get('data')
        obj = ProofModel.objects.first()
        proof = ProofForm(instance=obj)
        context = {'form': ProofreadingForm(), 'raqam': data, 'text': obj, 'proof': proof, 'title':'Proofreading'}
        return render(request, 'proofreading.html', context)

    def post(self, request):

        form = ProofreadingForm(request.POST, request.FILES)
        
        if form.is_valid():
            done = form.save()
            if done:
                print(request.POST)
                messages.success(request, 'Ваш запрос успешно отправлен, мы скоро свяжемся с вами!')
                return redirect('requests:proofreading')
        return render(request, 'proofreading.html', {'form': form})
       


class Peer_review(View):
    def get(self, request):
        context = {'form': PeerReviewForm(), 'title':'Peer Review'}
        return render(request, 'peer_review.html', context)

    def post(self, request):
        form = PeerReviewForm(request.POST, request.FILES)
        
        if form.is_valid():
            done = form.save()
            if done:
                messages.success(request, 'Ваш запрос успешно отправлен, мы скоро свяжемся с вами!')
                return redirect('requests:peer_review')
        return render(request, 'peer_review.html', {'form': form})
        


class BAKView(View):
    def get(self, request):
        obj = BakModel.objects.first()
        bak = BakForm(instance=obj)
        form = BaksForm()

        return render(request, 'bak.html', {'form': form, 'bak': bak, 'text': obj, 'title':'BAK'})

    def post(self, request):
        form = BaksForm(request.POST)
        
        if form.is_valid():
            done = form.save(commit=False)
            done.save()
            if done:
                messages.success(request, 'Ваш запрос успешно отправлен, мы скоро свяжемся с вами!')
                return redirect('requests:bak')
        return render(request, 'bak.html', {'validated': 'validated', 'form': form})
       

# updated
class ScopusView(View):
    def get(self, request):
        form = ScopusesForm()

        obj = ScopusModel.objects.first()
        scopus = ScopusForm(instance=obj)
        return render(request, 'scopus.html', {'form': form, 'text': obj, 'scopus': scopus, 'title': 'Scopus'})

    def post(self, request):
        form = ScopusesForm(request.POST)
        if form.is_valid():
            done = form.save(commit=False)
            done.save()
            if done:
                messages.success(request, 'Ваш запрос успешно отправлен, мы скоро свяжемся с вами!')
                return redirect('requests:scopus')
                
        return render(request, 'scopus.html', {'validated': 'validated', 'form': form})
       

class GrantsView(View):
    def get(self, request):
        obj = GrantsModel.objects.first()
        grant = GrantForm(instance=obj)
        form = GrantsForm()

        return render(request, 'grants.html', {'form': form, 'text': obj, 'grant': grant, 'title':'Grants'})

    def post(self, request):
        form = GrantsForm(request.POST, request.FILES)
        
        if form.is_valid():
            done = form.save(commit=False)
            done.save()
            if done:
                messages.success(request, 'Ваш запрос успешно отправлен, мы скоро свяжемся с вами!')
                return redirect('requests:grants')
            
        return render(request, 'grants.html', {'validated': 'validated', 'form': form})
    

class PatentsView(View):
    def get(self, request):
        form = PatentsForm()

        obj = PatentModel.objects.first()
        patent = PatentForm(instance=obj)
        return render(request, 'patents.html', {'form': form, 'text': obj, 'patent': patent, 'title':'Patent'})

    def post(self, request):
        form = PatentsForm(request.POST, request.FILES)
        
        if form.is_valid():
            done = form.save(commit=False)
            done.save()
            if done:
                messages.success(request, 'Ваш запрос успешно отправлен, мы скоро свяжемся с вами!')
                return redirect('requests:patents')
        return render(request, 'patents.html', {'validated': 'validated', 'form': form})
    
import docx2txt
import textract

import pptx
from pptx import Presentation
import codecs


class TranslationView(View):
    def get(self, request):
        
        text = TranslationModel.objects.first()
        translations = TranslationsForm(instance=text)

        form = TranslationForm()
        return render(request, 'translation.html', {'form': form, 'text': text, 'translations': translations, 'title':'Translation'})
        # request.session['back-trans'] = 'requests:translation'    

    def post(self, request):
        cost = TranslationCostModel.objects.first()
        form = TranslationForm(request.POST, request.FILES)

        user = request.user
        language = request.POST.get('language')
        research_area = request.POST.get('research_area')
        comment = request.POST.get('comment')
        file = request.FILES.get('file')
        char_amount = 0


        done = Translation.objects.create( char_amount=0, language=language,
                                            research_area=research_area, comment=comment, file=file, 
                                            )

        
        
        if str(done.file).endswith('.docx') or str(done.file).endswith('.pttx'):
            
        
            if done:
                text = textract.process(f'media/{done.file}')
                text2 = codecs.escape_decode(str(text))[0].decode()[2:-1]
                t = (',').join(text2.split('\n'))
                k = ''
                for i in t:
                    if i != ',':
                        k += i
                print(k)
                char_amount = len(k)

            done.char_amount = char_amount
            done.save()
            done.refresh_from_db()    

            product = 'Translation'
            amount = char_amount * cost.cost
            payment_type = 3
            payment_status = 0
            delivery_status = 0
            email = 'example@mail.com'
            phone = '998906568899'
            file = f'files/{file}'

            created = OrderModel.objects.create( product=product, amount=amount, payment_type=payment_type,
            payment_status=payment_status, delivery_status=delivery_status, email=email, phone=phone, file=file )
            if created:
                request.session['charcount'] = char_amount
                return redirect('orders:translation-payment')
            

        else:
            messages.warning(request, 'Убедитесь, что загруженый файл в формате docx или pptx.')
            done.delete()
            return redirect('requests:translation')

        
        # valid_ext = ['docx', 'pptx', 'txt']
        
        # files = str(file)
        # if '.' in files:
        #     ext = files.split('.')
        #     if ext in valid_ext:
        #         done = Translation.objects.create(user=user, word_amount=word_amount, language=language,
        #                                     research_area=research_area, comment=comment, file=file, 
        #                                     )
        #         if done:

        #             if files.endswith('.txt'):
        #                 with open(f'media/files/{file}',  'r') as read_file:
        #                     count_char = (',').join(read_file.read().split('\n'))
        #                     char_amount = len(count_char)
        #             elif files.endswith('.docx'):
        #                 my_text = docx2txt.process(f'media/files/{file}')
        #                 char_amount = len(my_text)
        #             elif files.endswith('.pptx'):
        #                 prs = Presentation(f'media/files/{file}')

        #                 # text_runs will be populated with a list of strings,
        #                 # one for each text run in presentation
        #                 text_runs = []

        #                 for slide in prs.slides:
        #                     for shape in slide.shapes:
        #                         if not shape.has_text_frame:
        #                             continue
        #                         for paragraph in shape.text_frame.paragraphs:
        #                             for run in paragraph.runs:
        #                                 text_runs.append(run.text)
        #                 # print(len((',').join(text_runs)))
        #                 char_amount = len((',').join(text_runs))

        #             else:
        #                 with open(f'media/files/{file}',  'r') as read_file:
        #                     count_char = (',').join(read_file.read().split('\n'))
        #                     char_amount = len(count_char)

        #             product = 'Translation'
        #             amount = char_amount * cost.cost
        #             payment_type = 3
        #             payment_status = 1
        #             delivery_status = 1
        #             email = user.email
        #             phone = user.phone
        #             file = f'files/{file}'

        #             created = OrderModel.objects.create(user=user, product=product, amount=amount, payment_type=payment_type,
        #             payment_status=payment_status, delivery_status=delivery_status, email=email, phone=phone, file=file )
        #             if created:
        #                 request.session['charcount'] = char_amount
        #                 return redirect('orders:translation-payment')

        #     else:
        #         messages.warning(request, 'bad file format file extention should .pptx, docx or txt')
        #         return redirect('requests:translation')

        # product = 'Translation'
        # amount = char_amount * cost.cost
        # payment_type = 3
        # payment_status = 0
        # delivery_status = 0
        # email = user.email
        # phone = user.phone
        # file = f'files/{file}'

        # created = OrderModel.objects.create(user=user, product=product, amount=amount, payment_type=payment_type,
        # payment_status=payment_status, delivery_status=delivery_status, email=email, phone=phone, file=file )
        # if created:
        #     request.session['charcount'] = char_amount
        #     return redirect('orders:translation-payment')
        
        return render(request, 'translation.html', {'validated': 'validated', 'form': form})


def after_login(request):
    request.session['after-login'] = 'pages:language_editing'
    return redirect('users:login')

    
